""" マルチスレッド対応のディレクトリ・ファイルを操作するクラスを提供する。
        Note: FileDataクラス・サブクラスはNotificationパターンによる実装をしている。エラーハンドリングするにはerrorsプロパティを参照する。
    ##  OSによらずディレクトリやフォルダ、ファイルのパスを表記する際のパスセパレータは「/」で統一するものとする。
    ### 標準サポートは.txt,.json,.csv,.xmlデータのみ
    ### 拡張サポートは.yaml,.docx,.xlsx,.xls,.pptx,.pdf
    TODO: Office関連,PDF関連のFileDataクラスはwrite()メソッドが未実装
"""

# 標準ライブラリのインポート
import re,json,csv,os,shutil,datetime,time,platform,zipfile,tarfile,hashlib,copy,base64,itertools,pathlib,io
import xml.etree.ElementTree as ET

IMPORTERROR_MESSAGE = "※このPython環境では{}には対応していません。"
# ファイルの排他制御(OSごとにことなるため分岐せざるを得ない)
try:
    if platform.system() == "Windows":
        import msvcrt
    else:
        import fcntl
except Exception as e:
    print(e)

try:
    import PIL

except Exception as e:
    print(e)
    print(IMPORTERROR_MESSAGE.format("画像データ"))

# 外部ライブラリのインポート(コメントはPython3.8.5)
# YAML対応する場合
# pip install ruamel.yaml
# ruamel.yaml      0.18.10
# ruamel.yaml.clib 0.2.12
try:
    import ruamel.yaml
except Exception as e:
    # print(e)
    print(IMPORTERROR_MESSAGE.format("YAMLデータ"))

# Office Word 対応する場合
# pip install python-docx
# python-docx 1.1.0
# - lxml      5.1.0
# - Pillow    10.2.0
try:
    import docx
except Exception as e:
    # print(e)
    print(IMPORTERROR_MESSAGE.format("Office Wordデータ"))

# Office Excel(.xlsx) 対応する場合
# pip install openpyxl
# openpyxl   3.1.2
# et-xmlfile 1.1.0
try:
    import openpyxl
except Exception as e:
    # print(e)
    print(IMPORTERROR_MESSAGE.format("Office Excel（.xlsx）データ"))

# Office Excel(.xls) 対応する場合
# pip install xlrd
# xlrd 2.0.1
try:
    import xlrd
except Exception as e:
    # print(e)
    print(IMPORTERROR_MESSAGE.format("Office Excel（.xls）データ"))

# Office PowerPoint(.pptx) 対応する場合
# pip install python-pptx
# python-pptx 1.0.2
# lxml        5.1.0
# Pillow      10.2.0
# XlsxWriter  3.1.9
try:
    import pptx
except Exception as e:
    # print(e)
    print(IMPORTERROR_MESSAGE.format("Office PowerPoint（.pptx）データ"))

# PDF対応する場合
# pip install PyMuPDF
# PyMuPDF  1.23.26
# PyMuPDFb 1.23.26
try:
    import fitz
except Exception as e:
    # print(e)
    print(IMPORTERROR_MESSAGE.format("PDFデータ"))

# HTML 対応する場合
# lxml        5.1.0

try:
    import lxml, lxml.html, lxml.etree
except Exception as e:
    # print(e)
    print(IMPORTERROR_MESSAGE.format("HTMLデータ"))

class DirData():
    """
    フォルダ・ディレクトリのデータを操作するクラス
    ```
        # 例)
        dd = DirData("etc/sample")
        print(dd.files)             # ファイルリストの参照
        dd.zip("sample.zip")        # 参照中のフォルダ・ディレクトリを圧縮
        print(dd.dir_name)          # ディレクトリパスからディレクトリ名を抽出
        dd.backup()                 # フォルダ・ディレクトリのバックアップ作成
        dd.delete()                 # フォルダ・ディレクトリの削除
        dd.rollback()               # フォルダ・ディレクトリのロールバック（バックアップがないとNG）
    ```
    """

    BACKUP_EXTEND = ".back"
    TAIL_BACKUP = "_back"
    TAIL_COPY = "_copy"
    _LINUX_EXTEND = ".tar.gz"
    _WINDOWS_EXTEND = ".zip"
    _DEFAULT_NAME = "tmp"
    _PLATFORM_NAME = platform.system()

    def __init__(self,dir_path):
        self.dir_path = dir_path
        if DirData._PLATFORM_NAME == "Windows":
            self.sep = DirData._WINDOWS_EXTEND
        else:
            self.sep = DirData._LINUX_EXTEND

    @property
    def dir_name(self):
        dn = ""
        dir_elms = self.dir_path.split("/")
        if len(dir_elms)>0:
            dn = f"{dir_elms[-1]}"
        else:
            dn = f"{DirData._DEFAULT_NAME}"
        return dn
    
    @property
    def dirs(self):
        tmp = {}
        for root, dir_s, files in os.walk(self.dir_path):
            for d in dir_s:
                dir_path = os.path.join(root,d)
                arcname = os.path.relpath(dir_path,start=self.dir_path)
                tmp[dir_path.replace("\\","/")] = arcname.replace("\\","/")
        return tmp
    
    @property
    def files(self):
        tmp = {}
        for root, dirs, files in os.walk(self.dir_path):
            for file in files:
                file_path = os.path.join(root,file)
                arcname = os.path.relpath(file_path,start=self.dir_path)
                tmp[file_path.replace("\\","/")] = arcname.replace("\\","/")
        return tmp
    
    @property
    def filenames(self):
        tmp = []
        for root, dirs, files in os.walk(self.dir_path):
            for file in files:
                tmp.append(file)
        return tmp

    def zip(self,zip_path=""):
        if zip_path == "":
            zip_path = "{}{}".format(self.dir_name,self.sep)
        with zipfile.ZipFile(zip_path,"w",zipfile.ZIP_DEFLATED) as zipf:
            for root, dirs, files in os.walk(self.dir_path):
                for file in files:
                    file_path = os.path.join(root,file)
                    arcname = os.path.relpath(file_path,start=self.dir_path)
                    zipf.write(file_path,arcname)
        return zip_path
    
    def targz(self,tar_path=""):
        if tar_path == "":
            tar_path = "{}{}".format(self.dir_name,self.sep)
        with tarfile.open(tar_path,"w:tar") as tar:
            tar.add(self.dir_path,arcname=os.path.basename(self.dir_path))
        return tar_path
    
    def copy(self,copy_path="",dirs_exist_ok=True):
        if copy_path == "" or copy_path == None:
            copy_path = f"{self.dir_path}{DirData.TAIL_COPY}"
        shutil.copytree(self.dir_path, copy_path,dirs_exist_ok=dirs_exist_ok)

    def backup(self):
        self.copy(f"{self.dir_path}{DirData.TAIL_BACKUP}")

    def rollback(self):
        dd = DirData(f"{self.dir_path}{DirData.TAIL_BACKUP}")
        dd.copy(self.dir_path)
        dd.delete()
    
    def delete(self):
        shutil.rmtree(self.dir_path)

    def create(self):
        os.makedirs(self.dir_path,exist_ok=True)


class FileData():
    """ ファイルの読み書きを行うクラス
        file_path:読み書きしたいjsonファイルパス
        data     :読み込んだファイルデータ
        read     :指定されたファイルを開き読み込んだ結果をdataへ格納する
        write    :指定されたファイルにdataを書き込む関数
        _read    :ファイルを読み込みをする関数（オーバーライド必須）
        _write   :ファイル書き込みをする関数（オーバーライド必須）
    """
    # 内部定数定義 """ 上書き可能 """
    BACKUP_EXTEND = ".back"             # バックアップを生成する際の拡張子
    SHA256_EXTEND = ".hash"             # SHA-256のハッシュ値ファイルの拡張子
    # デフォルトのエラーメッセージ（別モジュールから初期化上書きしてよい）
    ERROR_EMPTY_FILEPATH = "ファイルパスが設定されていません。"
    ERROR_FILE_LOCK      = "ファイルの排他制御開始に失敗しました。"
    ERROR_FILE_UNLOCK    = "ファイルの排他制御終了に失敗しました。"
    ERROR_MAX_RETRY      = "ファイル操作要求の試行回数が上限に達しました。"
    # ファイル操作要求のリトライ回数
    RETRY_COUNT = 10

    # 読み取り専用 """ 上書き禁止 """
    _PLATFORM_NAME = platform.system()  # Windows、Linuxなどの判定結果

    def __init__(self,file_path,encoding="utf-8",newline=None,first_read=True):
        # file path
        self.file_path = file_path # ファイルパス
        self.file_name = ""        # ファイル名
        self.abs_path  = ""        # ファイルの絶対パス
        self.suffix    = pathlib.Path(self.file_path).suffix # 拡張子

        # stack error messages
        self.messages = []         # エラーメッセージ
        self._data = None          # 読み込んだデータ
        self._backup_path = ""     # バックアップのパス
        self.encoding = encoding   # エンコード文字列設定値
        self.newline = newline     # 改行時の文字列
        if first_read == True:
            self.read()            # 初回読み込み 主にdata,filename の初期化
        else:
            pass                   # 初回読み込みをしなかった場合の処理 特になし

    @property
    def mode(self):
        return ""

    def read(self):
        """
        読み込み成功：True
        読み込み失敗：False
        """
        try:
            kwargs = {"encoding":self.encoding, "newline":self.newline}
            with open(self.file_path, "r{}".format(self.mode), **(kwargs if "b" not in self.mode else {})) as file:
                # self._data = json.load(file) などを行う
                self._read(file)
                self.file_name = file.name
                self.file_abs_path = os.path.abspath(file.name) if os.path.exists(file.name) else file.name
                return True
        except Exception as e:
            self.errors = e
            print(e)
        return False

    def _read(self,file):
        """
        fileを読み込みself_dataにデータを格納する
        """
        self.data = file.read()
        return None

    def write(self):
        """
        書き込み成功：True
        書き込み失敗：False
        """
        cnt = 0
        while True:
            if not self.is_locked:
                kwargs = {"encoding":self.encoding, "newline":self.newline}
                with open(self.file_path, "w{}".format(self.mode), **(kwargs if "b" not in self.mode else {})) as file:
                    try:
                        # json.dump(self._data,file,indent=4,ensure_ascii=False) などを行う
                        self.lock(file)
                        self._write(file)
                        return True
                    except FileNotFoundError as e:
                        self.errors = e
                        print(e)
                    except Exception as e:
                        self.errors = e
                        print(e)
                        return False
                    finally:
                        self.unlock(file)
            else:
                if cnt < FileData.RETRY_COUNT:
                    time.sleep(0.1)
                    cnt += 1
                else:
                    self.errors = FileData.ERROR_MAX_RETRY
                    return False
            
        
    def _write(self,file):
        """
        fileにself_dataを書き込む
        """
        file.write(self.data)
        return None
    
    def copy(self,copy_path):
        """
        コピー成功：True
        コピー失敗：False
        """
        try:
            self._copy(copy_path)
            return True
        except Exception as e:
            self.errors = e
            print(e)
            return False

    def _copy(self,copy_path):
        """
        file_path
        """
        tmp_path = ""
        if self.file_path == copy_path:
            tmp_path = ".".join([copy_path,self.now])
        else:
            tmp_path = copy_path
        shutil.copy(self.file_path,tmp_path)

    def delete(self):
        """
        削除成功：True
        削除失敗：False
        """
        cnt = 0
        while True:
            if self.is_locked == False:
                try:
                    self._delete()
                    return True
                except Exception as e:
                    self.errors = e
                    print(e)
                    return False
            else:
                if cnt < FileData.RETRY_COUNT:
                    time.sleep(0.1)
                    cnt += 1
                else:
                    self.errors = FileData.ERROR_MAX_RETRY
                    return False

        
    def _delete(self):
        """
        自ファイルを削除
        """
        os.remove(self.file_path)

    def backup(self):
        """
        バックアップ成功：True
        バックアップ失敗：False
        """
        try:
            self._backup()
            return True
        except Exception as e:
            self.errors = e
            print(e)
            return False

    def _backup(self):
        """ 
        バックアップファイル作成処理
        """
        self.copy(self.backup_path)

    def rollback(self):
        """
        復元成功：True
        復元失敗：False
        """
        cnt = 0
        while True:
            if self.is_locked == False:
                try:           
                        self._rollback()
                        return True
                except Exception as e:
                    self.errors = e
                    print(e)
                    return False
            else:
                if cnt < FileData.RETRY_COUNT:
                    time.sleep(0.1)
                    cnt += 1
                else:
                    self.errors = FileData.ERROR_MAX_RETRY
                    return False
        
    def _rollback(self):
        """
        ロールバック（復元処理）
        """
        result = False
        if self.file_path != "":
            try:
                fd = self.__class__(self.backup_path) #自インスタンスを新たに生成
                if fd.read():
                    fd.data
                    self.data = fd.data
                    if self.write():
                        result = True
                    else:
                        result = False
                else:
                    result = False
                return result
            except Exception as e:
                self.errors = e
                print(e)
                return result
        else:
            self.errors = FileData.ERROR_EMPTY_FILEPATH
            return result

    def hash(self):
        try:
            fd = TextData(f"{self.file_path}{FileData.SHA256_EXTEND}")
            fd.data = self.sha256
            print(fd.data)
            return fd.write()
            
        except Exception as e:
            print(e)
            self.errors = e
            return False    

    @property
    def backup_path(self):
        """ バックアップ先ファイルパスのゲッター
        """
        if self._backup_path == "":
            return self.file_path + FileData.BACKUP_EXTEND
        else:
            return self._backup_path

    @backup_path.setter
    def backup_path(self,bp):
        """ バックアップ先ファイルパスのセッター
        """
        self._backup_path = bp

    @property
    def data(self):
        """ ファイル内のデータのセッター（オーバーライド用）
        """
        return self._data
  
    @data.setter
    def data(self,data):
        """ ファイル内データのゲッター（オーバーライド用）
        """
        self._data = data

    @property
    def sha256(self):
        """ certutil -hashfile .\common\file.py
        
        """
        hash_sha256 = hashlib.sha256()
        with open(self.file_path,"rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_sha256.update(chunk)
        return hash_sha256.hexdigest()

    @property
    def now(self):
        """ タイムスタンプ生成関数（ファイル名重複回避用）
        """
        return datetime.datetime.now().strftime("%Y%m%d%H%M%S")
    @property
    def custom_timestamp(self):
        """ カスタムタイムスタンプ生成関数（ファイル名重複回避用）
        """
        return datetime.datetime.now().strftime("%Y/%m/%d/ %H:%M:%S")

    @property
    def is_error(self):
        """ エラーが発生しているかどうか
        """
        return len(self.messages)>0
    
    @property
    def errors(self):
        """ エラーメッセージのゲッター
        """
        return self.messages
    
    @errors.setter
    def errors(self,e):
        """ エラーメッセージのセッター
        """
        self.messages.append(e)

    def poperror(self):
        """ エラーメッセージを取得および返却し、オブジェクトが一時的に保持しているメッセージを破棄する
        """
        e = self.error[:]
        self.error = []
        return e
    
    @property
    def is_locked(self):
        try:
            with open(self.file_path, "r+") as f:
                if FileData._PLATFORM_NAME == "Windows":
                    msvcrt.locking(f.fileno(), msvcrt.LK_NBLCK,1)
                    msvcrt.locking(f.fileno(), msvcrt.LK_UNLCK,1)
                    return False
                else:
                    fcntl.flock(f.fileno(), fcntl.LOCK_EX | fcntl.LOCK_NB)
                    fcntl.flock(f.fileno(), fcntl.LOCK_UN)
                    return False
        except FileNotFoundError as e:
            # ファイルが存在しない＝ロックされていないでよい
            return False    
        except OSError as e:
            self.errors = e
            print(e)
        except BlockingIOError as e:
            self.errors = e
            print(e)
        except Exception as e:
            self.errors = e
            print(e)
        return True

    def lock(self,f):
        """ 排他制御開始
        """
        try:
            if FileData._PLATFORM_NAME == "Windows":
                msvcrt.locking(f.fileno(), msvcrt.LK_NBLCK,1)
            else:
                fcntl.flock(f.fileno(), fcntl.LOCK_EX | fcntl.LOCK_NB)
            return True
        except Exception as e:
            self.errors = e
            print(e)
            return False

    def unlock(self,f):
        """ 排他制御終了
        """
        try:
            if FileData._PLATFORM_NAME == "Windows":
                msvcrt.locking(f.fileno(), msvcrt.LK_UNLCK,1)
            else:
                fcntl.flock(f.fileno(), fcntl.LOCK_UN)
            return True
        except Exception as e:
            self.errors = e
            print(e)
            return False

class BinaryData(FileData):
    @property
    def mode(self):
        return "b"

class ImageController():

    def bytes_to_image(self,data: bytes) -> PIL.Image.Image:
        """ 
        DB / API / ファイルから読んだ画像を操作する場合に呼び出す

        :param data: 画像データ
        :return: PIL.imageデータ
        """
        return PIL.Image.open(io.BytesIO(data))
    
    def image_to_bytes(self,img: PIL.Image.Image, fmt="PNG") -> bytes:
        """ 画像をBase64変換の前段階のデータへ変換する
        """
        buf = io.BytesIO()
        img.save(buf, format=fmt)
        return buf.getvalue()
    
    def resize(
        self,
        img: PIL.Image.Image,
        size: tuple[int, int],
        keep_aspect=True
    ) -> PIL.Image.Image:
        if keep_aspect:
            img = img.copy()
            img.thumbnail(size)
            return img
        return img.resize(size)
    
    def normalize_mode(self,img: PIL.Image.Image, mode="RGB") -> PIL.Image.Image:
        """
            JPEGなどはRGBAなどのAlpha値がないため、エラーになる。そのため、この関数でRGBなどに変換する。
        """
        if img.mode != mode:
            return img.convert(mode)
        return img

    def image_to_data_url(
        self,
        img: PIL.Image.Image,
        fmt="PNG"
    ) -> str:
        """
        画像データをdata:urlに変換した文字列を取得する
        
        :param img: 画像データ
        :type img: PIL.Image.Image
        :param fmt: 画像のフォーマット
        :return: 画像データをdata:urlに変換したの文字列
        :rtype: str
        """
        raw = self.image_to_bytes(img, fmt)
        b64 = base64.b64encode(raw).decode("ascii")
        mime = f"image/{fmt.lower()}"
        return f"data:{mime};base64,{b64}"
    
    def data_url_to_image(self,data_url: str) -> PIL.Image.Image:
        """
        data:urlから画像データへ変換し画像を取得する
        
        :param data_url: data:url文字列
        :type data_url: str
        :return: 画像データ
        :rtype: Any
        """
        header, b64 = data_url.split(",", 1)
        data = base64.b64decode(b64)
        return self.bytes_to_image(data)
    
    def save(self,img: PIL.Image.Image, file_path:str, fmt=None):
        img.save(file_path, format=fmt)
        

class ImageData(BinaryData):
    def __init__(self,file_path,first_read=True):
        super().__init__(file_path=file_path,first_read=first_read)
        self.controller = ImageController()
        self.image = None
        if first_read:
            self.image = self.bytes_to_image()

    def bytes_to_image(self) -> PIL.Image.Image:
        """ DB / API / ファイルから読んだ画像を操作する場合に呼び出す
        :return: PIL.imageデータ
        """
        return self.controller.bytes_to_image(self.data)
    
    def image_to_bytes(self,fmt="PNG") -> bytes:
        """ 画像をBase64変換の前段階のデータへ変換する
        """
        return self.controller.image_to_bytes(self.image,fmt)
    
    def resize(self,size: tuple[int, int],keep_aspect=True) -> PIL.Image.Image:
        return self.controller.resize(self.image,size,keep_aspect)
    
    def normalize_mode(self, mode="RGB") -> PIL.Image.Image:
        """ JPEGなどはRGBAなどのAlpha値がないため、エラーになる。そのため、この関数でRGBなどに変換する。
        """
        return self.controller.normalize_mode(self.image,mode)

    def image_to_data_url(
        self,
        fmt="PNG"
    ) -> str:
        """
        画像データをdata:urlに変換した文字列を取得する
        
        :param fmt: 画像のフォーマット
        :return: 画像データをdata:urlに変換したの文字列
        :rtype: str
        """
        return self.controller.image_to_data_url(self.image,fmt)
    
    def data_url_to_image(self,data_url: str) -> PIL.Image.Image:
        """
        data:urlから画像データへ変換し画像を取得する
        
        :param data_url: data:url文字列
        :type data_url: str
        :return: 画像データ
        :rtype: Any
        """
        self.controller.data_url_to_image(data_url)
    
    def save(self,file_path:str, fmt=None):
        if self.image:
            self.controller.save(self.image,file_path,fmt)
        else:
            bd = BinaryData(file_path=file_path,first_read=False)
            bd.data = self.data
            if not (bd.data and bd.write()):
                self.errors = "{}書き込みに失敗しました。".format(file_path)

class JsonData(FileData):
    def __init__(self,file_path,encoding="utf-8",newline=None,first_read=True):
        super().__init__(file_path=file_path,encoding=encoding,newline=newline,first_read=first_read)

    def _read(self,file):
        self.data = json.load(file)

    def _write(self,file):
        json.dump(self.data,file,indent=4,ensure_ascii=False)

class TextData(FileData):
    def __init__(self,file_path,encoding="utf-8",newline="\n",first_read=True):
        super().__init__(file_path,encoding,newline,first_read)

    def _read(self, file):
        self.data = file.read()

    def _write(self, file):
        file.write(self.data)

class TextLineData(FileData):
    def __init__(self,file_path,encoding="utf-8",newline="",first_read=True):
        super().__init__(file_path,encoding,newline,first_read)

    def _read(self, file):
        self.data = file.readlines()

    def _write(self, file):
        linedata = []
        for d in self.data:
            if not d.endswith(self.newline):
                d = d + self.newline
            linedata.append(d)
        file.writelines(linedata)

class CSVData(FileData):
    def __init__(self,file_path,separator=",",encoding="utf-8-sig",first_read=True):
        self.sepa = separator
        self.encoding = encoding
        super().__init__(file_path,encoding,"",first_read)

    def _read(self,file):
        reader = csv.reader(file, delimiter=self.sepa)
        self.data = [row for row in reader]

    def _write(self,file):
        writer = csv.writer(file)
        writer.writerows(self.data)

    @property
    def head(self):
        return self.data[0] if self.data is not None else []
    
    @head.setter
    def head(self,arr):
        if self.data is None:
            self.data = [arr]
        else:
            if len(self.data)>0:
                self.data[0] = arr
            else:
                self.data.append(arr)
    
    @property
    def body(self):
        return self.data[1:] if self.data is not None and len(self.data) > 1 else []
    
    @body.setter
    def body(self,arr):
        hd = copy.deepcopy(self.head)
        bd = []
        tmp = [hd]
        for a in arr:
            tmp.append(a)
        self.data = tmp


class CSVDictData(FileData):
    def __init__(self,file_path,separator=",",encoding="utf-8-sig",first_read=True):
        self.sepa = separator
        self.encoding = encoding
        self.csv = CSVData(file_path,separator,encoding,first_read)
        super().__init__(file_path,encoding,"",first_read)

    def _read(self,file):
        reader = csv.DictReader(file,delimiter=self.sepa)
        self.head = reader.fieldnames
        self.body = [row for row in reader]
    
    def _write(self,file):
        file.newline = ""
        writer = csv.DictWriter(file,fieldnames=self.head,delimiter=self.sepa)
        # for row in self.data:
        #     writer.writerow(row)
        writer.writeheader()
        writer.writerows(self.body)


    @property
    def data(self):
        return self.csv.data
    
    @data.setter
    def data(self,data):
        self.csv.data = data

    @property
    def head(self):
        return self.csv.head
    
    @head.setter
    def head(self,arr):
        self.csv.head = arr

    @property
    def body(self):
        return self.csv.body
    
    @body.setter
    def body(self,arr):
        self.csv.body = arr

class XMLData(FileData):
    def __init__(self,file_path,encoding="utf-8",newline=None,first_read=True):
        self.tree = None
        super().__init__(file_path,encoding,newline)
    
    def _read(self,file):
        self.tree = ET.parse(self.file_path)
        self.data = self.tree.getroot()

    def _write(self,file):
        # ET.indent(self.tree,space="  ") #Python 3.9以上から使える
        self.tree.write(self.file_path,encoding=self.encoding,xml_declaration=True)

    def findall(self,name):
        return self.data.findall(name)

    def find(self,name,key,val):
        items = self.findall(name)
        target = None
        for item in items:
            if item.get(key) == val:
                target = item
        return target

    def create(self,itemName,attributes={},text=None):
        elm = ET.Element(itemName,attributes)
        if text is not None:
            elm.text = text
        # self.data.append(elm)
        return elm

class YamlData(FileData):
    def __init__(self,file_path,first_read=True):
        self.yaml = ruamel.yaml.YAML()
        self.file = None
        self._map = None
        super().__init__(file_path=file_path,encoding="utf-8",newline=None,first_read=first_read)

    def _read(self,file):
        self.file = file
        self.map = self.yaml.load(file)
        self.data = self.map

    def _write(self,file):
        self.unlock(file)
        # YAMLライブラリに扱わせるため、解除
        if self.map is None:
            self.yaml.dump(self.data,file)
        else:
            self.yaml.dump(self.map, file)
        self.lock(file)

    @property
    def data(self):
        return self.map if self.map is not None else self._data
    
    @data.setter
    def data(self,d):
        self._data = d
    
    @property
    def map(self):
        return self._map

    @map.setter
    def map(self,dic):
        self._map = dic

    def set_comment(self,key,val):
        if self.data.ca.items.get(key):
            self.data.ca.items[key][2].value = "# {}".format(val)
        else:
            self.data.yaml_add_eol_comment(val,key=key)


    def get_comments(self,dct=None):
        items = dct
        if dct is None:
            items = self.data
        comments = {}
        if items is not None:
            for k,v in items.items():
                comm = self.get_comment(k,items)
                if isinstance(v, dict):
                    comments[k] = {"comment":comm,"dict":self.get_comments(v)}
                else:
                    comments[k] = comm
        return comments
    
    def get_comment(self,key_or_index,container=None):
        """CommentedMap や CommentedSeq からコメントを取得"""
        if container is None:
            container = self.data
        if container is not None:
            if isinstance(container, ruamel.yaml.CommentedMap):
                if container.ca.items.get(key_or_index):
                    return container.ca.items[key_or_index][2].value.strip()
            elif isinstance(container, ruamel.yaml.CommentedSeq):
                # index番目の要素にコメントがあるか確認
                if container.ca.items.get(key_or_index):
                    return container.ca.items[key_or_index][2].value.strip()
        return None
    
    def array(self,data):
        arr = []
        for k,v in data.items():
            arr.append({k:v})
        return arr

    def dictionary(self,data):
        dic = {}
        for d in data:
            for k,v in d.items():
                dic[k] = v
        return dic
    
    def update(self,newdata):
        if isinstance(newdata,(dict,tuple)):
            for k,v in newdata.items():
                self.data[k] = v
        elif isinstance(newdata,list):
            cnt = 0
            for l in newdata:
                self.data[cnt] = l

class WordData(FileData):
    def __init__(self,file_path,first_read=True):
        self.document = docx.Document(file_path)
        self.lines = []
        self.images = []
        super().__init__(file_path=file_path,encoding="utf-8",newline=None,first_read=first_read)
    
    def _read(self,file):
        self.lines = []
        paragraphs = []
        for i,p in enumerate(self.document.paragraphs):
            info = {"text":"","images":[]}
            images = info["images"]
            for run_idx, run in enumerate(p.runs):
                drawing_elems = run._element.findall(".//a:blip", namespaces={
                    # Word（Office Open XML）で使われる XML 名前空間（namespace）の識別子
                    "a": "http://schemas.openxmlformats.org/drawingml/2006/main" 
                })
                for blip in drawing_elems:
                    rId = blip.get(docx.oxml.ns.qn("r:embed"))
                    image_part = self.document.part.related_parts[rId]
                    image_bytes = image_part.blob
                    images.append({
                        "para_idx":i,      #段落ID
                        "run_idx":run_idx, #段落内の文書ID
                        "data":image_bytes,
                        "name":image_part.partname.replace("/","-"),
                        "suffix":image_part.content_type.split("/")[-1],
                    })

            if p.text.strip():
                self.lines.append(p.text.strip())
                info["text"] = p.text.strip()
            paragraphs.append(info)

        self.data = paragraphs

    def _write(self,file):
        raise Exception("書き込み処理は未実装です。")
    
    def text(self):
        return "\n".join([l for l in self.lines])

    def output_images(self):
        for d in self.data:
            images = d["images"]
            for img in images:
                fp = "{}-{}-{}-{}.{}".format(self.file_path,img.get("para_idx",""),img.get("run_idx",""),img.get("name","img.png").split(".")[0],img.get("suffix",""))
                imgdata = ImageData(fp,first_read=False)
                imgdata.data = img.get("data",None)
                if not (imgdata and imgdata.write()):
                    self.errors = "{}:書き込みに失敗しました。\n{}".format(fp,"\n".join(imgdata.errors))

class ExcelData(FileData):
    def __init__(self,file_path,first_read=True,has_header=False,skip_empty_rows=True):
        self.has_header = has_header
        self.skip_empty_rows = skip_empty_rows
        self.document = None
        super().__init__(file_path=file_path,encoding="utf-8",newline=None,first_read=first_read)

    def _read(self,file):
        if self.suffix == ".xls":
            self.data = self.load_xls()
        elif self.suffix == ".xlsx" or self.suffix == ".xlsm":
            self.data = self.load_xlsx()
        else:
            raise Exception("サポートしていないデータです。")
        
    @property
    def sheet_names(self):
        return self.data.keys() if self.data is not None else []
    
    def col(self,sheet_name,index):
        return [d[index] for d in self.data[sheet_name]["body"]] if self.data is not None else None
    
    def row(self,sheet_name,index):
        return self.data[sheet_name]["body"][index] if self.data is not None else None

    def load_xls(self):
        wb = xlrd.open_workbook(self.file_path,formatting_info=False)
        self.document = wb
        result = {}
        for si in range(wb.nsheets):
            ws = wb.sheet_by_index(si)
            rows = []
            org_headers = []
            if ws.nrows == 0:
                result[ws.name] = {"header":org_headers,"body":rows,"images":[]}
                continue
            if self.has_header:
                org_headers = ws.row_values(0)
                headers = [str(h).strip() if h not in {None, ""} else f"col_{i+1}" for i, h in enumerate(org_headers)]
                for r in range(1, ws.nrows):
                    row_vals = ws.row_values(r)
                    if self.skip_empty_rows and all((v == "" or v is None) for v in row_vals):
                        continue
                    if len(row_vals) < len(headers):
                        row_vals += [None] * (len(headers) - len(row_vals))
                    elif len(row_vals) > len(headers):
                        row_vals = row_vals[:len(headers)]

                    rows.append({h: v for h, v in zip(headers, row_vals)})


                result[ws.name] = {"header":org_headers,"body":rows,"images":[]}

            else:
                for r in range(ws.nrows):
                    row_vals = ws.row_values(r)
                    if self.skip_empty_rows and all((v == "" or v is None) for v in row_vals):
                        continue
                    rows.append(list(row_vals))
                result[ws.name] = {"header":org_headers,"body":rows,"images":[]}
        return result

    def load_xlsx(self):
        wb = openpyxl.load_workbook(self.file_path,data_only=True,read_only=False)
        self.document = wb
        result = {}
        for ws in wb.worksheets:

            # 画像データの取得
            images = []
            for i, img in enumerate(ws._images):
                pil_img = img._data()
                suffix = ""
                if isinstance(pil_img, bytes):
                    suffix = PIL.Image.open(io.BytesIO(pil_img)).format.lower()
                else:
                    suffix = pil_img.format.lower()

                # 画像のアンカーポイント（厳密な座標までは取得ロジックが構成できていない）
                point = {"r":img.anchor._from.row + 1,"c":img.anchor._from.col + 1}
                images.append({
                    "index": i + 1,
                    "data": pil_img,
                    "suffix": suffix,
                    "name": f"{ws.title}_{i}.{suffix}",
                    "point":point,
                })


            rows_iter = ws.iter_rows(values_only=True)
            rows = []
            # ヘッダーあり
            if self.has_header:
                # 1行目をヘッダーとして採用
                try:
                    headers = next(rows_iter)
                except StopIteration:
                    result[ws.title] = {"header":[],"body":[],"images":images}
                    continue
                org_headers = list(headers) if headers is not None else []
                headers = [str(h).strip() if h is not None else f"col_{i+1}" for i, h in enumerate(headers)]
                for row in rows_iter:
                    if self.skip_empty_rows and (row is None or all(cell is None for cell in row)):
                        continue
                    # 行の長さがヘッダーと異なる場合のケア
                    row = list(row) if row is not None else []
                    row = ["" if v is None else v for v in row]
                    if len(row) < len(headers):
                        row += [None] * (len(headers) - len(row))
                    elif len(row) > len(headers):
                        row = row[:len(headers)]
                    rows.append({h: v for h, v in zip(headers, row)})
            
            # ヘッダーなし
            else:
                org_headers = []
                for row in rows_iter:
                    if self.skip_empty_rows and (row is None or all(cell is None for cell in row)):
                        continue
                    
                    row = list(row) if row is not None else []
                    row = ["" if v is None else v for v in row]
                    rows.append(row)
            
                
            result[ws.title] = {"header":org_headers,"body":rows,"images":images}
        wb.close()
        return result

    def _write(self,file):
        raise Exception("書き込み処理は未実装です。")
    
    def output_images(self):
        for k in self.data.keys():
            info = self.data[k]
            images = info["images"]
            for img in images:
                fp = "{}-{}-{}.{}".format(self.file_path,img.get("index",""),img.get("name","img.png").split(".")[0],img.get("suffix",""))
                imgdata = ImageData(fp,first_read=False)
                imgdata.data = img["data"]
                if not(imgdata.data and imgdata.write()):
                    self.errors = "{}書き込みに失敗しました。".format(fp)

    
class PowerPointData(FileData):
    def __init__(self,file_path,first_read=True):
        self.document = pptx.Presentation(file_path)
        self.images = []
        super().__init__(file_path=file_path,encoding="utf-8",newline=None,first_read=first_read)
    
    def _read(self,file):
        presentations = []
        texts = []
        images = []
        for i, slide in enumerate(self.document.slides, start=1):
            slide_info = {"index":i,"shapes":[]}
            slide_texts = []
            for shape in self.group_shapes(slide.shapes):
                if hasattr(shape,"text") and shape.text.strip():
                    t = shape.text.strip()
                else:
                    t = ""
                slide_info["shapes"].append({"text":t,"type":shape.shape_type,"name":shape.name,"object":shape})
                slide_texts.append(t)
                if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE: # MSO_SHAPE_TYPE.PICTURE 13
                    images.append({
                        "name":shape.image.filename,
                        "data":shape.image.blob,
                        "slide-index":i,
                    })

            presentations.append(slide_info)
            if slide_texts:
                texts.append(f"# Slide {i}\n" + "\n".join(slide_texts))

        self.data = presentations
        self.pages = presentations
        self.text = "\n\n".join(texts)
        self.images = images


    def _write(self,file):
        raise Exception("書き込み処理は未実装です。")
    
    def group_shapes(self,shapes):
        for shape in shapes:
            yield shape
            if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.GROUP:
                yield from shapes(shape.shapes)

    def output_images(self):
        cnt = 1
        for shape in self.images:
            fp = "{}-{}-{}-{}".format(self.file_path,shape.get("slide-index",""),cnt,shape.get("name","img"))
            cnt += 1
            img = ImageData(fp,first_read=False)
            img.data = shape.get("data",None)
            if not(img.data and img.write()):
                self.errors = "{}-{}書き込みに失敗しました。\n{}".format(shape.get("name",""),shape.get("slide-index",""),"\n".join(img.errors))
                

class PDFData(FileData):

    # 1ポイント: 1/72インチ
    DEFAULT_WIDTH  = 595 
    DEFAULT_HEIGHT = 842 
    DEFAULT_MERGIN = 50

    def __init__(self,file_path,first_read=True):
        self.pages = []
        self.document = None
        self.width  = PDFData.DEFAULT_WIDTH
        self.height = PDFData.DEFAULT_HEIGHT
        self.mergin = PDFData.DEFAULT_MERGIN
        super().__init__(file_path=file_path,encoding="utf-8",newline=None,first_read=first_read)
    
    def _read(self,file):
        self.document = fitz.open(self.file_path)
        parts = []
        for i, page in enumerate(self.document, start=1):
            text = page.get_text("text")
            images = page.get_images(full=True)
            info ={"text":text.strip() if text else "","images":[]}
            for img in images:
                xref = img[0] # xref index
                image_name = img[7] # name index
                base_image = self.document.extract_image(xref)
                image_bytes  = base_image["image"]
                image_suffix = base_image["ext"]
                info["images"].append({"data":image_bytes,"name":image_name,"suffix":image_suffix,"page":i,"info":img})

            parts.append(info)

        self.pages = parts
        self.data = parts
        self.document.close()

    def _write(self,file):
        """ オーバーライド用
            Note:PDF制御は複雑であるためここでは、実装しない。サブクラスで詳細なメソッドとともに実装するべき。
        """
        raise Exception("書き込み処理は未実装です。")
    

    @property
    def text(self):
        return "\n\n".join([p.get("text","") for p in self.pages])

    def make(self, file_path, texts=[]):
        """ 新規作成でPDFファイルを生成します。
        :param file_path: 保存先ファイルパス
        :param texts: ページあたりのテキストデータの配列
        """

        doc = fitz.open()
        self.document = doc
        
        for idx, txt in enumerate(texts):
            # 新しいページを追加
            page = doc.new_page(width=self.width, height=self.height)

            # 文字列を書き込む (左上から)
            rect = fitz.Rect(self.mergin, self.mergin, self.width - self.mergin, self.height - self.mergin)  # マージン 50pt
            page.insert_textbox(
                rect,
                txt,
                fontsize=12,
                fontname="helv",
                align=fitz.TEXT_ALIGN_LEFT,
                color=(0, 0, 0)  # 黒色
            )

        # 1ページ以上のテキストがあれば保存
        if len(texts) > 0:
            # ファイルを保存
            doc.save(file_path)
        doc.close()

    def output_images(self):
        pages = self.data
        cnt = 1
        for page in pages:
            images = page["images"]
            for img in images:
                fp = "{}-{}-{}-{}.{}".format(self.file_path,img.get("page",""),cnt,img.get("name","img"),img.get("suffix",""))
                imgdata = ImageData(fp,first_read=False)
                imgdata.data = img.get("data",None)
                if not (imgdata.data and imgdata.write()):
                    self.errors = "{}書き込みに失敗しました。".format(fp)
                    cnt += 1


class HTMLData(FileData):

    def __init__(self,file_path,encoding="utf-8",newline=None,first_read=True):
        self.document = None
        self.targets = None  # Xpathで取得した対象のDOM .text = "", .set("key","value")
        super().__init__(file_path=file_path,encoding=encoding,newline=newline,first_read=first_read)
    
    def _read(self,file):
        super()._read(file)
        self.document = lxml.html.fromstring(self.data)

    def _write(self,file):
        self.data = lxml.etree.toString(self.document,pretty_print=True,encoding="unicode")
        return super()._write(file)
    
    def xpath(self,x):
        """
        HTMLからXPATHによるDOMの取得をする        
            Args: param x: Xpath表現文字列 "//div[@class='item']"
            Return: Array: items
        """
        self.targets = self.document.xpath(x)
        return self.targets
    
    @property
    def body(self):
        return self.xpath("//body")[0]
    
    @property
    def title(self):
        return self.xpath("//title/text()")[0].strip()
    
    @property
    def fullText(self):
        for script in self.body.xpath("//script | //style | //noscript"):
            script.drop_tree()
        t = "".join(self.body.itertext())
        return re.sub(r"\s+", " ", t).strip()
    
    def analize(self,data):
        self.data = data
        self.document = lxml.html.fromstring(self.data)

    def append(self,dom=None,target=None):
        if target:
            target.append(dom)
        else:
            self.body.append(dom)

    def remove(self,x):
        items = self.document.xpath(x)
        for item in items:
            item.getparent().remove(item)

    def create(self,name,options={}):
        dom = lxml.etree.Element(name)
        t = options.get("text",None)
        if t:
            dom.text = t
        for k,v in options:
            dom.set(k,v)

        return dom


class Crypt():

    DECRYPT_ERROR = "複合化に失敗しました。"
    ENCRYPT_ERROR = "暗号化に失敗しました。"

    def __init__(self,key="",encoding="utf-8",decoding="ascii"):
        self.key = key
        self.encoding = encoding
        self.decoding = decoding
        self.token = ""
        self.byte = ""

    def encrypte(self,text=""):
        """ 暗号化:UTF-8文字列を鍵でXORし、Base84で可読化する
        text: 暗号化するデータ
        """
        data = text.encode(self.encoding)
        k = self.key.encode(self.encoding)
        xored = bytes([b ^ k for b, k in zip(data, itertools.cycle(k))])
        self.byte = base64.urlsafe_b64encode(xored).decode(self.decoding)
        return self.byte

    def decrypte(self,token):
        """ 複合化:Base64文字列をデコードし、同じ鍵でXOR複合する
        token: 複合化するデータ（暗号化されたデータ）
        """
        xored = base64.urlsafe_b64decode(token.encode(self.decoding))
        k = self.key.encode(self.encoding)
        data = bytes([b ^ k for b, k in zip(xored, itertools.cycle(k))])
        self.token = data.decode(self.encoding)
        return self.token

class CryptFileData(FileData):
    ERROR_MAKE_CRYPT = "暗号化ファイルの生成に失敗しました。"
    ERROR_NOT_FOUND_KEY = "複合化キーがありません。"
    """
    Base64/XORで暗号化されたファイルの生成・読み込みクラス
    """
    def __init__(self,file_path="",key_path="",encoding="utf-8",decoding="ascii",newline=None,first_read=True):
        self.key_path = key_path
        self.key = FileData(key_path,encoding,newline,True)
        self.decoding = decoding
        self.crypt = Crypt(self.key.data,encoding,decoding)
        self.txt = ""
        super().__init__(file_path=file_path,encoding=encoding,newline=newline,first_read=first_read)

        if self.key.data:
            self.errors = CryptFileData.ERROR_NOT_FOUND_KEY


    def encrypte(self,txt=""):
        """ 暗号化メソッド
        Args:
            txt string: 暗号化する文字列
        Return: String | None
        """
        try:
            return self.crypt.encrypte(txt)
        except Exception as e:
            print(e)
            print(Crypt.ENCRYPT_ERROR)
            return None
    
    def decrypte(self):
        """ 複合化メソッド
        Args:-
        Return: String | None
        """
        try:
            return self.crypt.decrypte(self.txt)
        except Exception as e:
            print(e)
            print(Crypt.DECRYPT_ERROR)
            return None

    @property    
    def decrypted(self):
        return self.crypt.token
    
    @property        
    def encrypted(self):
        return self.crypt.byte
    
    def _read(self, file):
        """
        ファイルを読み込みself.txtにテキストデータを格納。複合化したデータをself.dataに格納。
        """
        self.txt = file.read()
        self.data = self.decrypte()
        

    def _write(self, file):
        """
        ファイルにself.dataを暗号化したデータを書き込む。
        """
        file.write(self.encrypte(self.data))

    def make(self,txt=""):
        fd = FileData(file_path=self.file_path,encoding=self.encoding,newline=self.newline,first_read=False)
        fd.data = self.encrypte(txt)
        result = fd.write()
        if result:
            self.read()
        else:
            self.errors = CryptFileData.ERROR_MAKE_CRYPT
        return result
    
    def remake(self):
        """
        現在の複合化したデータを基に、ファイルを作り直す。
        """
        fd = FileData(self.file_path,encoding=self.encoding,newline=self.newline,first_read=False)
        fd.data = self.data
        result = fd.write()
        if result == False:
            self.errors = CryptFileData.ERROR_MAKE_CRYPT
        return result

class CryptTextData(CryptFileData):
    """
    Base64/XORで暗号化されたテキストファイルの生成・読み込みクラス
    ```
        # 使用例
        ctd = CryptJsonData("暗号化されたファイルパス","鍵ファイルパス","utf-8","ascii",None,True)
        print("複合化データ",ctd.data)
        print("暗号化書き込み",ctd.write())
    ```
    """
    def __init__(self,file_path="",key_path="",encoding="utf-8",decoding="ascii",newline=None,first_read=True):
        super().__init__(file_path=file_path,key_path=key_path,encoding=encoding,decoding=decoding,newline=newline,first_read=first_read)

    def _read(self, file):
        return super()._read(file)

    def _write(self, file):
        return super()._write(file)
    
    def make(self,txt=""):
        """
        暗号化されたファイルを生成するメソッド        
        ```
        ctd = CryptTextData("etc/test.txt","etc/key.txt")
        ctd.make('秘匿メッセージ')
        print(ctd.data)
        ```
        """
        return super().make(txt)
    
    def remake(self):
        return super().remake()


class CryptJsonData(CryptFileData):
    """
    Base64/XORで暗号化されたJSONファイルの生成・読み込みクラス
    ```
        # 使用例
        cjd = CryptJsonData("暗号化されたファイルパス","鍵ファイルパス","utf-8","ascii",None,True)
        print("複合化データ",cjd.data)
        print("暗号化書き込み",cjd.write())
    ```
    """
    def __init__(self,file_path="",key_path="",encoding="utf-8",decoding="ascii",newline=None,first_read=True):
        super().__init__(file_path=file_path,key_path=key_path,encoding=encoding,decoding=decoding,newline=newline,first_read=first_read)

    def _read(self, file):
        self.txt = file.read()
        self.data = json.loads(self.decrypte())
        
    def _write(self, file):
        file.write(self.encrypte(json.dumps(self.data)))

    def make(self,data={}):
        """
        暗号化されたファイルを生成するメソッド   
        ```
        cjd = CryptJsonData("etc/test.json","etc/key.txt","utf-8","ascii",None,False)
        cjd.make({
            "pass":"秘密",
            "path":{
                "AAA":"aaa",
                "BBB":"bbb",
                "CCC":"ccc"
            },
        })
        print(cjd.data)
        ```
        """
        jd = JsonData(self.file_path,encoding=self.encoding,newline=self.newline,first_read=False)
        jd.data = self.encrypte(json.dumps(data))
        result = jd.write()
        if result:
            self.read()
        else:
            self.errors = CryptFileData.ERROR_MAKE_CRYPT
        return result
    
    def remake(self):
        jd = JsonData(self.file_path,encoding=self.encoding,newline=self.newline,first_read=False)
        jd.data = self.data
        result = jd.write()
        if result == False:
            self.errors = CryptFileData.ERROR_MAKE_CRYPT
        return result
    
if __name__ == "__main__":
    """
    py -m common.file
    """

    # htmld = HTMLData("etc/data/sample.html")
    # print(htmld.title)
    # print(htmld.fullText)

    dd = DirData("etc")
    print(dd.dir_name)
    print(dd.dirs)

    # cjd = JsonData("etc/test.json")
    # cjd = CryptJsonData("etc/test.json","etc/key.txt","utf-8","ascii",None,True)
    # print("複合化データ",cjd.data)
    # cjd.remake()

    # bfd = ImageData("etc/data/sample.png")
    # print(bfd.image)
    # print(bfd.image_to_data_url("PNG"))

    # wfd = WordData("etc/data/sample.docx")
    # print(wfd.data)
    # wfd.output_images()

    # efd = ExcelData("etc/data/sample.xlsx",has_header=True)
    # for c in efd.col("Sheet1",1):
    #     print(c)
    # print(efd.sheet_names)
    # print(efd.data["Sheet3"])
    # efd.output_images()

    # pdffd = PDFData("etc/data/sample.pdf")
    # print(pdffd.data)

    # ppfd = PowerPointData("etc/data/sample.pptx")
    # # print(ppfd.data)
    # # print(ppfd.images)
    # ppfd.output_images()
    # print(ppfd.errors)

    # pdffd = PDFData("etc/data/sample.pdf")
    # print(pdffd.text)
    # pdffd.output_images()
    # print(pdffd.errors)


    # ppffd2 = PDFData("etc/data/sample2.pdf",first_read=False)
    # ppffd2.make(ppffd2.file_path,["aaaa\nThis is sample. This is sample. This is sample. This is sample.This is sample. This is sample. This is sample. This is sample.","bbbbb"])

    ## 以下使用例

    # def test1():
    #     """ TextData write()テスト"""
    #     d = TextData("etc/sample/sample.txt","utf-8",None,False)
    #     d.data = datetime.datetime.strftime(datetime.datetime.now(),"%Y%M%D %H%m%S")
    #     print("data",d.data)
    #     print("result",d.write())
    #     print("ファイルの書き込みを確認してください。")

    # def test2():
    #     """ TextData read()テスト"""
    #     d = TextData("etc/sample/sample.txt","utf-8",None,False)
    #     print("before",d.data)
    #     print("result",d.read())
    #     print("after",d.data)

    # def test3():
    #     """ TextData copy()テスト"""
    #     d = TextData("etc/sample/sample.txt")
    #     print("result",d.copy("etc/sample/sample.copy.txt"))

    # def test4():
    #     """ TextData backup()テスト"""
    #     d = TextData("etc/sample/sample.txt")
    #     print("result",d.backup())

    # def test5():
    #     """ TextData delete()テスト"""
    #     d = TextData("etc/sample/sample.txt")
    #     print("result",d.delete())

    # def test6():
    #     """ TextData rollbak()テスト"""
    #     d = TextData("etc/sample/sample.txt",first_read=False)
    #     print("result",d.rollback())

    # def test11():
    #     """ JsonData write()テスト"""
    #     d = JsonData("etc/sample/sample.json",first_read=False)
    #     d.data = {"now":datetime.datetime.strftime(datetime.datetime.now(),"%Y%M%D %H%m%S")}
    #     print("data",d.data)
    #     print("result",d.write())
    #     print("ファイルの書き込みを確認してください。")

    # def test12():
    #     """ JsonData read()テスト"""
    #     d = JsonData("etc/sample/sample.json","utf-8",None,False)
    #     print("before",d.data)
    #     print("result",d.read())
    #     print("after",d.data)

    # def test13():
    #     """ JsonData copy()テスト"""
    #     d = JsonData("etc/sample/sample.json")
    #     print("result",d.copy("etc/sample/sample.copy.json"))

    # def test14():
    #     """ JsonData backup()テスト"""
    #     d = JsonData("etc/sample/sample.json")
    #     print("result",d.backup())

    # def test15():
    #     """ JsonData delete()テスト"""
    #     d = JsonData("etc/sample/sample.json")
    #     print("result",d.delete())

    # def test16():
    #     """ JsonData rollbak()テスト"""
    #     d = JsonData("etc/sample/sample.json",first_read=False)
    #     print("result",d.rollback())

    # def test21():
    #     """ CSVData write()テスト"""
    #     d = CSVData("etc/sample/sample.csv",first_read=False)
    #     d.data = [
    #         {1,datetime.datetime.strftime(datetime.datetime.now(),"%Y%M%D %H%m%S")},
    #         {2,datetime.datetime.strftime(datetime.datetime.now(),"%Y%M%D %H%m%S")},
    #     ]
    #     print("data",d.data)
    #     print("result",d.write())
    #     print("ファイルの書き込みを確認してください。")

    # def test22():
    #     """ CSVData read()テスト"""
    #     d = CSVData("etc/sample/sample.csv","utf-8",None,False)
    #     print("before",d.data)
    #     print("result",d.read())
    #     print("after",d.data)

    # def test23():
    #     """ CSVData copy()テスト"""
    #     d = CSVData("etc/sample/sample.csv")
    #     print("result",d.copy("etc/sample/sample.copy.csv"))

    # def test24():
    #     """ CSVData backup()テスト"""
    #     d = CSVData("etc/sample/sample.csv")
    #     print("result",d.backup())

    # def test25():
    #     """ CSVData delete()テスト"""
    #     d = CSVData("etc/sample/sample.csv")
    #     print("result",d.delete())

    # def test26():
    #     """ CSVData rollbak()テスト"""
    #     d = CSVData("etc/sample/sample.csv",first_read=False)
    #     print("result",d.rollback())

    # def test31():
    #     """ CSVDictData write()テスト"""
    #     d = CSVDictData("etc/sample/sample.dict.csv",first_read=False)
    #     d.head = ["row","time"]
    #     print(d.head)
    #     d.body = [
    #         {"row":1,"time":datetime.datetime.strftime(datetime.datetime.now(),"%Y%M%D %H%m%S")},
    #         {"row":2,"time":datetime.datetime.strftime(datetime.datetime.now(),"%Y%M%D %H%m%S")},
    #     ]
    #     print(d.body)
    #     print("data",d.data)
    #     print("result",d.write())
    #     print("ファイルの書き込みを確認してください。")

    # def test32():
    #     """ CSVDictData read()テスト"""
    #     d = CSVDictData("etc/sample/sample.dict.csv",first_read=False)
    #     print("before",d.data)
    #     print("result",d.read())
    #     print("head",d.head)
    #     print("body",d.body)
    #     print("after",d.data)

    # def test33():
    #     """ CSVDictData copy()テスト"""
    #     d = CSVDictData("etc/sample/sample.dict.csv")
    #     print("result",d.copy("etc/sample/sample.dict.copy.csv"))

    # def test34():
    #     """ CSVDictData backup()テスト"""
    #     d = CSVDictData("etc/sample/sample.dict.csv")
    #     print("result",d.backup())

    # def test35():
    #     """ CSVDictData delete()テスト"""
    #     d = CSVDictData("etc/sample/sample.dict.csv")
    #     print("result",d.delete())

    # def test36():
    #     """ CSVDictData rollbak()テスト"""
    #     d = CSVDictData("etc/sample/sample.dict.csv",first_read=False)
    #     print("result",d.rollback())


    # def test41():
    #     """ YamlData write()テスト"""
    #     d = YamlData("etc/sample/sample.yaml",first_read=True)
    #     d.data = {
    #         "row1":"This is sample.",
    #         "row2":{"row":1,"time":datetime.datetime.strftime(datetime.datetime.now(),"%Y%M%D %H%m%S")},
    #         "row3":{"row":2,"time":datetime.datetime.strftime(datetime.datetime.now(),"%Y%M%D %H%m%S")},
    #     }
    #     print("result",d.write())
    #     print("data",d.data)
    #     print("ファイルの書き込みを確認してください。")

    # def test42():
    #     """ YamlData read()テスト"""
    #     d = YamlData("etc/sample/sample.yaml",first_read=False)
    #     print("before",d.data)
    #     print("result",d.read())
    #     print("after",d.data)

    # def test43():
    #     """ YamlData copy()テスト"""
    #     d = YamlData("etc/sample/sample.yaml")
    #     print("result",d.copy("etc/sample/sample.copy.yaml"))

    # def test44():
    #     """ YamlData backup()テスト"""
    #     d = YamlData("etc/sample/sample.yaml")
    #     print("result",d.backup())

    # def test45():
    #     """ YamlData delete()テスト"""
    #     d = YamlData("etc/sample/sample.yaml")
    #     print("result",d.delete())

    # def test46():
    #     """ YamlData rollbak()テスト"""
    #     d = YamlData("etc/sample/sample.yaml",first_read=False)
    #     print("result",d.rollback())

    # def test47():
    #     """ YamlData comment()テスト"""
    #     d = YamlData("etc/sample/sample.yaml")
    #     print("data",d.data,d.map)
    #     print("comment",d.get_comment("row1"))
    #     print("comments",d.get_comments())
    #     d.set_comment("row2","bbbb")
    #     print("result",d.write())
    #     print(d.get_comment("row2"))


    # def test51():
    #     """ TextLineData write()テスト"""
    #     d = TextLineData("etc/sample/sample.line.txt",newline="\n",first_read=True)
    #     d.data = ["aaa","bbb","cccc"]
    #     print("result",d.write())
    #     print("data",d.data)
    #     print("ファイルの書き込みを確認してください。")

    # def test52():
    #     """ TextLineData read()テスト"""
    #     d = TextLineData("etc/sample/sample.line.txt",newline="\n",first_read=False)
    #     print("before",d.data)
    #     print("result",d.read())
    #     print("after",d.data)

    # def test53():
    #     """ TextLineData copy()テスト"""
    #     d = TextLineData("etc/sample/sample.line.txt")
    #     print("result",d.copy("etc/sample/sample.copy.txt"))

    # def test54():
    #     """ TextLineData backup()テスト"""
    #     d = TextLineData("etc/sample/sample.line.txt")
    #     print("result",d.backup())

    # def test55():
    #     """ TextLineData delete()テスト"""
    #     d = TextLineData("etc/sample/sample.line.txt")
    #     print("result",d.delete())

    # def test56():
    #     """ TextLineData rollbak()テスト"""
    #     d = TextLineData("etc/sample/sample.line.txt",first_read=False)
    #     print("result",d.rollback())



    # test_list = [
    #     # FileDataのサブクラスの動作確認
    #     # TextData
    #     test1,
    #     test2,
    #     test3,
    #     test4,
    #     test5,
    #     test6,
    #     # JSONData
    #     test11,
    #     test12,
    #     test13,
    #     test14,
    #     test15,
    #     test16,
    #     # CSVData
    #     test21,
    #     test22,
    #     test23,
    #     test24,
    #     test25,
    #     test26,
    #     # CSVDictData
    #     test31,
    #     test32,
    #     test33,
    #     test34,
    #     test35,
    #     test36,
    #     # YAMLData
    #     test41,
    #     test42,
    #     test43,
    #     test44,
    #     test45,
    #     test46,
    #     test47,
    #     # TextLineData
    #     test51,
    #     test52,
    #     test53,
    #     test54,
    #     test55,
    #     test56,

    # ]
    # for t in test_list:
    #     print(t.__name__,t.__doc__)
    #     t()

    pass

